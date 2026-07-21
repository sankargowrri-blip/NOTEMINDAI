"""Export service: PDF, DOCX, TXT, HTML, Markdown, PPTX."""
import io
import markdown as md_lib
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt


def export_to_pdf(title: str, text: str) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(title, styles["Title"]))
    story.append(Spacer(1, 0.5 * cm))
    for line in text.splitlines():
        if line.startswith("# "):
            story.append(Paragraph(line[2:], styles["Heading1"]))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles["Heading2"]))
        elif line.startswith("### "):
            story.append(Paragraph(line[4:], styles["Heading3"]))
        elif line.startswith("```") or line.startswith("    "):
            story.append(Preformatted(line, styles["Code"]))
        elif line.strip():
            story.append(Paragraph(line, styles["BodyText"]))
        else:
            story.append(Spacer(1, 0.3 * cm))
    doc.build(story)
    return buffer.getvalue()


def export_to_docx(title: str, text: str) -> bytes:
    doc = Document()
    doc.add_heading(title, level=0)
    for line in text.splitlines():
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.strip():
            doc.add_paragraph(line)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def export_to_txt(text: str) -> bytes:
    return text.encode("utf-8")


def export_to_html(title: str, text: str) -> bytes:
    html_body = md_lib.markdown(text, extensions=["tables", "fenced_code", "nl2br"])
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>{title}</title>
<style>
  body {{ font-family: Georgia, serif; max-width: 900px; margin: 40px auto; padding: 0 20px; line-height: 1.7; }}
  h1, h2, h3 {{ color: #1a1a2e; }}
  code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 4px; }}
  pre {{ background: #f4f4f4; padding: 16px; border-radius: 8px; overflow-x: auto; }}
  table {{ border-collapse: collapse; width: 100%; }}
  td, th {{ border: 1px solid #ddd; padding: 8px; }}
</style>
</head>
<body>
<h1>{title}</h1>
{html_body}
</body>
</html>"""
    return html.encode("utf-8")


def export_to_markdown(title: str, text: str) -> bytes:
    content = f"# {title}\n\n{text}"
    return content.encode("utf-8")


def export_to_pptx(title: str, text: str) -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    # Split content into slides (one per heading or every 10 lines)
    lines = text.splitlines()
    current_title = "Overview"
    current_content = []

    def _add_slide(t, c):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = t
        tf = slide.placeholders[1].text_frame
        tf.text = "\n".join(c[:15])  # max 15 lines per slide

    for line in lines:
        if line.startswith("# ") or line.startswith("## "):
            if current_content:
                _add_slide(current_title, current_content)
            current_title = line.lstrip("# ").strip()
            current_content = []
        elif line.strip():
            current_content.append(line.strip())
    if current_content:
        _add_slide(current_title, current_content)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


EXPORT_FORMATS = {
    "pdf": ("application/pdf", ".pdf", export_to_pdf),
    "docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx", export_to_docx),
    "txt": ("text/plain", ".txt", lambda title, text: export_to_txt(text)),
    "html": ("text/html", ".html", export_to_html),
    "md": ("text/markdown", ".md", export_to_markdown),
    "pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx", export_to_pptx),
}


def export_note(title: str, text: str, fmt: str) -> tuple[bytes, str, str]:
    """Returns (file_bytes, content_type, filename)."""
    if fmt not in EXPORT_FORMATS:
        raise ValueError(f"Unsupported format: {fmt}")
    content_type, ext, fn = EXPORT_FORMATS[fmt]
    data = fn(title, text)
    filename = f"{title.replace(' ', '_')}{ext}"
    return data, content_type, filename
